"""
pptx_loader.py
--------------
Extract text from PowerPoint (.pptx) files and convert to the document format
expected by build_index():  {paper_id, title, full_text}

Each slide's title + body text is extracted and concatenated with a
"--- Slide N ---" separator so the chunker can split naturally.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation


def _normalize_pdf_text(text: str) -> str:
    """
    Fix pypdf extraction artifacts: collapse multiple spaces and join
    single-word lines (word-wrap newlines) back into flowing paragraphs.
    Double newlines (paragraph / page breaks) are preserved.
    """
    import re
    # Collapse multiple spaces
    text = re.sub(r" {2,}", " ", text)
    # Join single newlines that are word-wrap artifacts.
    # Don't touch \n\n (paragraph breaks) or \n--- (page-separator headers).
    text = re.sub(r"\n(?!\n|---)", " ", text)
    # Collapse any new runs of multiple spaces created by the join
    text = re.sub(r" {2,}", " ", text)
    return text


def _load_pdf(file_path: Path) -> dict:
    """Extract text from a PDF and return a document dict."""
    import pypdf

    paper_id = file_path.stem
    parts = []
    title = paper_id
    with open(file_path, "rb") as f:
        reader = pypdf.PdfReader(f)
        for i, page in enumerate(reader.pages, start=1):
            raw = (page.extract_text() or "").strip()
            if not raw:
                continue
            # Extract title from the first non-empty line of raw text (before normalization)
            if title == paper_id:
                for line in raw.splitlines():
                    line = line.strip()
                    if line:
                        title = line
                        break
            text = _normalize_pdf_text(raw)
            parts.append(f"--- Page {i} ---\n{text}")

    full_text = "\n\n".join(parts)
    return {"paper_id": paper_id, "title": title, "full_text": full_text}


def _slide_text(slide) -> str:
    """Return all text from a single slide, title first."""
    parts = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            line = " ".join(run.text for run in para.runs).strip()
            if line:
                parts.append(line)
    return "\n".join(parts)


def load_pptx(file_path: str | Path) -> dict:
    """
    Load a single .pptx file and return a document dict.

    Returns
    -------
    {
        "paper_id":  str  — filename without extension
        "title":     str  — first slide title or filename
        "full_text": str  — all slide text concatenated
    }
    """
    file_path = Path(file_path)
    prs = Presentation(str(file_path))

    paper_id = file_path.stem  # filename without extension
    title = paper_id           # fallback if no title slide

    slide_blocks = []
    for i, slide in enumerate(prs.slides, start=1):
        text = _slide_text(slide).strip()
        if not text:
            continue

        # Use the first non-empty slide text as document title
        if i == 1 and title == paper_id:
            first_line = text.splitlines()[0]
            if first_line:
                title = first_line

        slide_blocks.append(f"--- Slide {i} ---\n{text}")

    full_text = "\n\n".join(slide_blocks)
    return {"paper_id": paper_id, "title": title, "full_text": full_text}


def load_pptx_dir(directory: str | Path) -> list[dict]:
    """
    Load all .pptx and .pdf files from a directory (skips Office temp files).

    Returns a list of document dicts, one per file.
    """
    directory = Path(directory)

    pptx_files = [f for f in sorted(directory.glob("*.pptx")) if not f.name.startswith("~$")]
    ppt_files  = [f for f in sorted(directory.glob("*.ppt"))  if not f.name.startswith("~$")]
    pdf_files  = sorted(directory.glob("*.pdf"))
    files = pptx_files + ppt_files + pdf_files

    if not files:
        print(f"[pptx_loader] No .pptx or .pdf files found in {directory}")
        return []

    documents = []
    for f in files:
        try:
            if f.suffix.lower() == ".pdf":
                doc = _load_pdf(f)
            else:
                doc = load_pptx(f)
            print(f"[pptx_loader] Loaded: {f.name!r} — {len(doc['full_text'])} chars, title={doc['title'][:60]!r}")
            documents.append(doc)
        except Exception as e:
            print(f"[pptx_loader] WARNING: Failed to load {f.name!r}: {e}")

    print(f"[pptx_loader] Total: {len(documents)} presentations loaded.")
    return documents
